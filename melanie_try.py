import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
from io import BytesIO
import json
from PIL import Image
from fpdf import FPDF
from datetime import datetime
from openai import OpenAI
from langchain.chains import ConversationChain
from langchain.memory import ConversationBufferMemory
from langchain.chat_models import ChatOpenAI
import dotenv
import os
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import csv
import io

# --- Initialize and Settings ---
dotenv.load_dotenv()

def initialize_client(api_key):
    """Initialize OpenAI client with the provided API key."""
    return OpenAI(api_key=api_key) if api_key else None

def execute_code_and_generate_image(code: str):
    """執行 Python 程式碼並生成圖表圖片"""
    try:
        exec(code, globals())
        if 'plt' in globals():
            buf = BytesIO()
            plt.savefig(buf, format="png")
            buf.seek(0)
            return buf
        return None
    except Exception as e:
        return str(e)

def read_file_content(file):
    """讀取上傳文件的內容"""
    try:
        if file.type == "application/pdf":
            reader = PdfReader(file)
            return "".join([page.extract_text() for page in reader.pages])
        elif file.type == "text/plain":
            return file.getvalue().decode("utf-8")
        elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            doc = Document(file)
            return "\n".join([para.text for para in doc.paragraphs])
        elif file.type == "text/csv":
            content = io.StringIO(file.getvalue().decode("utf-8"))
            return "\n".join([", ".join(row) for row in csv.reader(content)])
        elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            prs = Presentation(file)
            return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    except Exception as e:
        st.error(f"File processing failed: {e}")
    return None

def main():
    st.set_page_config(page_title="Chatbot + Data Analysis", page_icon="🤖", layout="centered")
    st.title("🤖 Chatbot + 📊 Data Analysis + 🧠 Memory")

    with st.sidebar:
        st.subheader("🔒 Enter Your API Key")
        api_key = st.text_input("OpenAI API Key", type="password")

        # 文件上傳功能
        st.subheader("📄 Upload a Document")
        uploaded_doc = st.file_uploader("Choose a file (txt, pdf, docx, csv, pptx):", 
                                       type=["txt", "pdf", "docx", "csv", "pptx"])
        if uploaded_doc:
            doc_content = read_file_content(uploaded_doc)
            if doc_content:
                st.write("### File Content:")
                st.text_area("Extracted Content:", doc_content, height=300)
                st.session_state.messages.append({"role": "user", "content": doc_content})

        # CSV 上傳
        st.subheader("📂 Upload a CSV File")
        uploaded_file = st.file_uploader("Choose a CSV file:", type=["csv"])
        csv_data = pd.read_csv(uploaded_file) if uploaded_file else None

        # 其他功能按鈕
        if st.sidebar.button("🗑️ Clear Memory"):
            st.session_state.memory.clear()
            st.session_state.messages = []
            st.success("Memory cleared!")
        if st.sidebar.button("💾 Save Conversation as JSON"):
            save_conversation_to_json()

    # 初始化聊天記錄與記憶體
    if "messages" not in st.session_state:
        st.session_state.messages = []
    if "conversation" not in st.session_state:
        if api_key:
            st.session_state.chat_model = ChatOpenAI(model="gpt-4-turbo", temperature=0.5, openai_api_key=api_key)
            st.session_state.memory = ConversationBufferMemory()
            st.session_state.conversation = ConversationChain(
                llm=st.session_state.chat_model, memory=st.session_state.memory)
        else:
            st.warning("⬅️ Please enter the API key to initialize the chatbot.")
            return

    # 聊天與回應顯示
    user_input = st.chat_input("Hi! Ask me anything...")
    if user_input:
        st.session_state.messages.append({"role": "user", "content": user_input})
        with st.chat_message("user"):
            st.write(user_input)

        with st.spinner("Thinking..."):
            try:
                response = st.session_state.conversation.run(user_input)
                with st.chat_message("assistant"):
                    st.write(response)

                # 檢測並執行程式碼
                if "```python" in response:
                    code_start = response.find("```python") + len("```python")
                    code_end = response.find("```", code_start)
                    code = response[code_start:code_end].strip()
                    image_buffer = execute_code_and_generate_image(code)
                    if isinstance(image_buffer, BytesIO):
                        st.image(image_buffer, caption="Generated Chart", use_container_width=True)
                    elif isinstance(image_buffer, str):
                        st.error(f"Code execution error: {image_buffer}")
            except Exception as e:
                st.error(f"Error: {e}")

if __name__ == "__main__":
    main()
